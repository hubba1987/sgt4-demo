from pptx import Presentation

def set_text_by_element_addr(elem_addr, text_value):
    global slide
    addr_split_list = elem_addr.split(":")
    shape_idx = int(addr_split_list[0])
    paragrapth_idx = int(addr_split_list[1])
    run_idx = int(addr_split_list[2])
    slide.shapes[shape_idx].text_frame.paragraphs[paragrapth_idx].runs[run_idx].text = text_value

prs = Presentation('CV_template.pptx')
#slide_layout = prs.slide_layouts[0]
slide = prs.slides[0]
ns=0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            np=0
            for paragraph in shape.text_frame.paragraphs:
                nr=0
                for run in paragraph.runs:
                    print(f"{ns}:{np}:{nr} {run.text}")
                    nr += 1
                np += 1
        else:
            print(shape)
        ns += 1
#slide.shapes[1].text_frame.paragraphs[0].text = "Aliaksei Yazhou"
set_text_by_element_addr("1:0:0", "Aliaksei Yazhou")
set_text_by_element_addr("2:0:0", "JA11")
set_text_by_element_addr("3:3:0", "In IT for 16 years already...")

prs.save("new.pptx")