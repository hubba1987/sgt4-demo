from datetime import datetime
from flask import Flask, render_template, request, redirect, url_for, send_from_directory, send_file
from pptx import Presentation
from pptx.text.text import _Paragraph
import copy, os
app = Flask(__name__)


    

@app.route('/')
def index():
   print('Request for index page received')
   return render_template('index.html')

@app.route('/favicon.ico')
def favicon():
    return send_from_directory(os.path.join(app.root_path, 'static'),
                               'favicon.ico', mimetype='image/vnd.microsoft.icon')

@app.route('/generate', methods=['GET'])
def generate():

    def set_text_by_element_addr(slide, elem_addr, text_value):
        #global slide
        addr_split_list = elem_addr.split(":")
        shape_idx = int(addr_split_list[0])
        paragrapth_idx = int(addr_split_list[1])
        run_idx = int(addr_split_list[2])
        slide.shapes[shape_idx].text_frame.paragraphs[paragrapth_idx].runs[run_idx].text = text_value

    def append_skills_paragraphs(shape, skill_level, skills_list):
        global heading_1_style, heading_2_style   
        for skill in skills_list:
            skill_p = shape.text_frame.add_paragraph()
            heading_run = skill_p.add_run()
            heading_run.text = skill

    def append_oneline_paragraph(shape, heading_style, text):
        level_p = shape.text_frame.add_paragraph()
        heading_run = level_p.add_run()
        heading_run._r.insert(0,heading_style)
        heading_run.text = text

    name = request.args.get('name')
    level = request.args.get('level')
    background = request.args.get('background')

    prs = Presentation('CV_template.pptx')
    slide = prs.slides[0]

    ### Manipulating avatar start
    new_picture = slide.shapes[7]._parent.add_picture("avatar.jpg", 485172, 162987, height = slide.shapes[7].height)
    new_picture.auto_shape_type = slide.shapes[7].auto_shape_type
    old_pic = slide.shapes[7]._element
    new_pic = new_picture._element
    old_pic.addnext(new_pic)  # method on lxml _Element, moves in this case, doesn't copy
    old_pic.getparent().remove(old_pic)
    ### Manipulating avatar end

    heading_1_style = copy.deepcopy(slide.shapes[4].text_frame.paragraphs[0].runs[0].font._rPr)
    heading_2_style = copy.deepcopy(slide.shapes[4].text_frame.paragraphs[1].runs[0].font._rPr)

    ### Clear skills and language block
    slide.shapes[4].text_frame.clear()
    par_to_remove = slide.shapes[4].text_frame.paragraphs[0]._element
    par_to_remove.getparent().remove(par_to_remove)

    append_oneline_paragraph(slide.shapes[4], heading_1_style, "Functional/technical skills")
    append_oneline_paragraph(slide.shapes[4], heading_2_style, "EXPERT")

    ### Add "Functional/technical skills" heading
    expert_skills = ('AWS', 'Jenkins', 'Terraform')
    append_skills_paragraphs(slide.shapes[4], "Expert", expert_skills)
    append_oneline_paragraph(slide.shapes[4], heading_2_style, " ")
    append_oneline_paragraph(slide.shapes[4], heading_2_style, "INTERMEDIATE")
    interm_skills = ('GCP', 'Gitab', 'Ansible', 'Linux', 'Python', 'Confluence', 'Helm')

    append_skills_paragraphs(slide.shapes[4], "Intermediate", interm_skills)

    begginer_skills = ('Java', 'MySQL', 'Kubernetes')
    append_oneline_paragraph(slide.shapes[4], heading_2_style, " ")
    append_oneline_paragraph(slide.shapes[4], heading_2_style, "BEGGINER")
    append_skills_paragraphs(slide.shapes[4], "Begginer", begginer_skills)

        #slide.shapes[4].text_frame.paragraphs = paragraph

    set_text_by_element_addr(slide, "1:0:0", name)
    set_text_by_element_addr(slide, "2:0:0", level)
    set_text_by_element_addr(slide, "3:3:0", background)

    prs.save("new.pptx")
    return send_file("new.pptx")


if __name__ == '__main__':
   app.run()